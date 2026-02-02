# 必要なライブラリ
# pip install python-pptx pillow

from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os

IMG_DIR = 'pdf_images'
PPTX_PATH = 'Integrated_Childcare_System_Foundation.pptx'

# スライドサイズ（A4横, 単位:EMU）
SLIDE_WIDTH = Inches(11.69)
SLIDE_HEIGHT = Inches(8.27)

def images_to_pptx(img_dir, pptx_path):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    import re
    def sort_key(filename):
        m = re.search(r'page_(\d+)\\.png|page_(\d+)\.png', filename)
        if m:
            return int(m.group(1) or m.group(2))
        return 0
    img_files = sorted([f for f in os.listdir(img_dir) if f.endswith('.png')], key=sort_key)
    for img_file in img_files:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 白紙
        img_path = os.path.join(img_dir, img_file)
        # 画像サイズ取得
        with Image.open(img_path) as im:
            img_width_px, img_height_px = im.size
        # 画像の物理サイズ（インチ）
        img_width_in = img_width_px / 96
        img_height_in = img_height_px / 96
        # スライドサイズ（インチ）
        slide_width_in = SLIDE_WIDTH.inches
        slide_height_in = SLIDE_HEIGHT.inches
        # スライドに収まる最大サイズを計算（アスペクト比維持）
        width_ratio = slide_width_in / img_width_in
        height_ratio = slide_height_in / img_height_in
        scale = min(width_ratio, height_ratio)
        disp_width = img_width_in * scale
        disp_height = img_height_in * scale
        # 中央配置
        left = Inches((slide_width_in - disp_width) / 2)
        top = Inches((slide_height_in - disp_height) / 2)
        slide.shapes.add_picture(img_path, left, top, width=Inches(disp_width), height=Inches(disp_height))
    prs.save(pptx_path)
    print(f'PPTX保存完了: {pptx_path}')

if __name__ == '__main__':
    images_to_pptx(IMG_DIR, PPTX_PATH)
